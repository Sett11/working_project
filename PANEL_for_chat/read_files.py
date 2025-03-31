import os
import io
from typing import Dict, List, Tuple, Union
import pdfplumber  # для работы с PDF
import docx  # python-docx для работы с DOCX
from pptx import Presentation  # python-pptx для работы с PPTX
from PIL import Image
from logs import log_event
from action_model import describe_images_sync
from config import MAX_CHARS


# start

def safe_extract_text(func):
    """Декоратор для безопасного извлечения текста с логированием ошибок"""
    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            log_event("ERROR", f"Failed to extract text: {str(e)}")
            return "", [], []
    return wrapper

@safe_extract_text
def extract_text_from_txt(file_path: str) -> Tuple[str, List[Image.Image]]:
    """Извлечение текста из TXT файлов"""
    log_event("TXT_PROCESS", f"Starting TXT processing for file: {file_path}")
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    log_event("TXT_PROCESS", f"TXT processing completed. Text length: {len(text)}")
    log_event("TXT_CONTENT", f"First 500 chars of extracted text: {text[:500]}")
    return text, []  # Возвращаем текст и пустой список изображений

@safe_extract_text
def extract_text_from_pdf(file_path: str) -> Tuple[str, List[Image.Image]]:
    """Извлечение текста и изображений из PDF"""
    text = ""
    images = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            # Извлекаем текст
            page_text = page.extract_text() or ""
            text += page_text
            # Извлекаем изображения
            for image in page.images:
                try:
                    # Получаем изображение как bytes
                    image_bytes = image["stream"].get_data()
                    log_event("INFO", f"Image size: {len(image_bytes) / (1024 * 1024):.2f} MB")
                    image = Image.open(io.BytesIO(image_bytes))
                    images.append(image)
                except Exception as e:
                    log_event("ERROR", f"Failed to process image from PDF: {str(e)}")
    return text, images

@safe_extract_text
def extract_text_from_docx(file_path: str) -> Tuple[str, List[Image.Image]]:
    """Извлечение текста и изображений из DOCX"""
    text = ""
    images = []
    log_event("DOCX_PROCESS", f"Starting DOCX processing for file: {file_path}")
    doc = docx.Document(file_path)
    log_event("DOCX_PROCESS", f"Document loaded successfully. Paragraphs: {len(doc.paragraphs)}, Tables: {len(doc.tables)}")
    # Извлекаем текст из параграфов
    for i, para in enumerate(doc.paragraphs):
        # Добавляем текст параграфа
        para_text = para.text.strip()
        if para_text:  # Логируем только непустые параграфы
            # log_event("DOCX_PARAGRAPH", f"Paragraph {i+1}: {para_text[:100]}...")
            text += para_text + "\n"
        # Проверяем стиль параграфа
        if para.style.name.startswith('Heading'):
            log_event("DOCX_HEADING", f"Heading {i+1}: {para_text[:100]}...")
            text += "\n"  # Добавляем дополнительный перенос для заголовков
    # Извлекаем текст из таблиц
    for table_idx, table in enumerate(doc.tables):
        log_event("DOCX_TABLE", f"Processing table {table_idx+1} with {len(table.rows)} rows")
        for row_idx, row in enumerate(table.rows):
            row_text = ""
            for cell_idx, cell in enumerate(row.cells):
                # Добавляем текст из ячейки
                cell_text = cell.text.strip()
                if cell_text:  # Логируем только непустые ячейки
                    log_event("DOCX_CELL", f"Table {table_idx+1}, Row {row_idx+1}, Cell {cell_idx+1}: {cell_text[:100]}...")
                    row_text += cell_text + " "
                # Проверяем параграфы внутри ячейки
                for para in cell.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        row_text += para_text + " "
            if row_text.strip():
                text += row_text.strip() + "\n"
        text += "\n"  # Добавляем перенос между таблицами
    # Извлекаем изображения
    image_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                image = Image.open(io.BytesIO(image_data))
                images.append(image)
                image_count += 1
                log_event("INFO", f"Image size: {len(image_data) / (1024 * 1024):.2f} MB")
            except Exception as e:
                log_event("ERROR", f"Failed to process image from DOCX: {str(e)}")
    # Очищаем текст от лишних пробелов и переносов
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    log_event("DOCX_PROCESS", f"DOCX processing completed. Text length: {len(text)}, Images: {image_count}")
    log_event("DOCX_CONTENT", f"First 500 chars of extracted text: {text[:500]}")
    return text, images

@safe_extract_text
def extract_text_from_pptx(file_path: str) -> Tuple[str, List[Image.Image]]:
    """Извлечение текста и изображений из PPTX"""
    text = ""
    images = []
    log_event("PPTX_PROCESS", f"Starting PPTX processing for file: {file_path}")
    prs = Presentation(file_path)
    log_event("PPTX_PROCESS", f"Presentation loaded successfully. Total slides: {len(prs.slides)}")
    for slide_idx, slide in enumerate(prs.slides):
        log_event("PPTX_SLIDE", f"Processing slide {slide_idx + 1}")
        slide_text = []
        # Обрабатываем заголовок слайда
        try:
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    log_event("PPTX_TITLE", f"Slide {slide_idx + 1} title: {title_text[:100]}...")
                    slide_text.append(f"Заголовок: {title_text}")
        except Exception as e:
            log_event("WARNING", f"Could not process title for slide {slide_idx + 1}: {str(e)}")
        # Обрабатываем все фигуры на слайде
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # Обрабатываем текст из фигур
                if hasattr(shape, "text") and shape.text.strip():
                    shape_text = shape.text.strip()
                    log_event("PPTX_SHAPE", f"Slide {slide_idx + 1}, Shape {shape_idx + 1}: {shape_text[:100]}...")
                    slide_text.append(shape_text)
                # Обрабатываем таблицы
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    log_event("PPTX_TABLE", f"Slide {slide_idx + 1}, Table {shape_idx + 1}: {len(table.rows)} rows")
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))                
                # Обрабатываем изображения
                if hasattr(shape, "shape_type") and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        image = Image.open(image_stream)
                        images.append(image)
                        log_event("INFO", f"Image size: {len(image.tobytes()) / (1024 * 1024):.2f} MB")
                        log_event("PPTX_IMAGE", f"Successfully extracted image from slide {slide_idx + 1}")
                    except Exception as e:
                        log_event("ERROR", f"Failed to process image from PPTX slide {slide_idx + 1}: {str(e)}")
            except Exception as e:
                log_event("ERROR", f"Error processing shape {shape_idx + 1} on slide {slide_idx + 1}: {str(e)}")
        # Обрабатываем заметки к слайду
        try:
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    log_event("PPTX_NOTES", f"Slide {slide_idx + 1} notes: {notes_text[:100]}...")
                    slide_text.append(f"Заметки: {notes_text}")
        except Exception as e:
            log_event("WARNING", f"Could not process notes for slide {slide_idx + 1}: {str(e)}")
        # Добавляем текст слайда в общий текст
        if slide_text:
            text += f"\nСлайд {slide_idx + 1}:\n" + "\n".join(slide_text) + "\n"
    # Очищаем текст от лишних пробелов и переносов
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    log_event("PPTX_PROCESS", f"PPTX processing completed. Text length: {len(text)}, Images: {len(images)}")
    log_event("PPTX_CONTENT", f"First 500 chars of extracted text: {text[:500]}")
    return text, images

def process_document(file_path: str) -> Dict[str, Union[str, bool]]:
    """Обрабатывает документ и извлекает из него текст и описания изображений"""
    file_name = os.path.basename(file_path)
    file_ext = os.path.splitext(file_name)[1].lower()  # Приводим расширение к нижнему регистру
    log_event("INFO", f"Processing document: {file_name} (extension: {file_ext})")
    # Словарь соответствия расширений и функций обработки
    processors = {
        '.txt': (extract_text_from_txt, []),
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.pptx': extract_text_from_pptx
    }
    try:
        if file_ext not in processors:
            log_event("WARNING", f"Неподдерживаемый формат файла: {file_ext}")
            return {
                "content": f"Неподдерживаемый формат файла: {file_ext}",
                "truncated": False
            }
        # Получаем функцию обработки и дополнительные аргументы
        processor, extra_args = processors[file_ext] if isinstance(processors[file_ext], tuple) else (processors[file_ext], [])
        # Обрабатываем документ
        text, images = processor(file_path, *extra_args)
        # Проверяем, что текст не пустой
        if not text.strip():
            log_event("WARNING", f"Документ {file_name} пуст или не содержит текста")
            return {
                "content": f"Документ {file_name} пуст или не содержит текста",
                "truncated": False
            }
        if file_ext == '.pptx':
            text = text.split('\nСлайд')[0]
        log_event("INFO", f"Документ {file_name} успешно обработан. Длина текста: {len(text)}, Количество изображений: {len(images)}")
        # Обрабатываем изображения и добавляет их описания к тексту
        new_text, descriptions = text, ''
        if images:
            max_weight_images = sorted(images, key=lambda x: x.size[0] * x.size[1])[-20:]
            descriptions = '\n\n'.join(f'<text_to_image>{i2t}</text_to_image>\n' for i2t in describe_images_sync(max_weight_images) if i2t != '')
            log_event("INFO", f"Описание изображений: {descriptions}")
            new_text= text + '\n\n' + descriptions
        # Проверяем длину и при необходимости обрезаем
        if len(new_text) > MAX_CHARS:
            new_text = text[:MAX_CHARS - len(descriptions) - 2] + '\n\n' + descriptions
            log_event("WARNING", f"Документ {file_name} был обрезан (превышает {MAX_CHARS} символов)")
            return {
                "content": new_text,
                "truncated": True
            }
        log_event("INFO", f"Длина конечного содержимого документа: {len(text)}")
        return {
            "content": new_text,
            "truncated": False
        }
    except Exception as e:
        log_event("ERROR", f"Ошибка при обработке документа {file_path}: {str(e)}")
        return {
            "content": f"Ошибка при обработке файла: {str(e)}",
            "truncated": False
        }